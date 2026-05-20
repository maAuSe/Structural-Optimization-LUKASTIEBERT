"""Build the MTOP presentation in the KU Leuven template."""
import os
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, "kul ppt template.pptx")
OUT = os.path.join(ROOT, "MTOP_presentation.pptx")
FIG = os.path.join(ROOT, "matlab", "figures")
RFIG = os.path.join(ROOT, "report", "figures")
GEN = os.path.join(ROOT, "presentation_build")

TEAL = RGBColor(0x1D, 0x8D, 0xB0)
SLATE = RGBColor(0x2F, 0x4D, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5C, 0x6E, 0x7A)

SW, SH = 12192000, 6858000
MARGIN = 576000
CONTENT_T = 1410000
CONTENT_B = 6150000
CONTENT_H = CONTENT_B - CONTENT_T
COLGAP = 320000
LEFT_X = MARGIN
LEFT_W = 4790000
RIGHT_X = LEFT_X + LEFT_W + COLGAP
RIGHT_W = SW - MARGIN - RIGHT_X
FOOTER = "Structural Optimization (B-KUL-H0T88A)"

prs = Presentation(TEMPLATE)
NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def nb(s):
    """Apply non-breaking spaces so dimensions / units never split across lines."""
    s = s.replace(" × ", " × ")
    s = s.replace(" → ", " → ")
    s = s.replace("K u = f", "K u = f")
    s = re.sub(r"(\d) ([s%])(?=[\s.,;:)]|$)", "\\1 \\2", s)
    return s


def layout(name):
    for m in prs.slide_masters:
        for lay in m.slide_layouts:
            if lay.name == name:
                return lay
    raise KeyError(name)


def delete_all_slides():
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        prs.part.drop_rel(sid.get(NS_R))
        lst.remove(sid)


def ph(slide, idx):
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def drop(shape):
    shape._element.getparent().remove(shape._element)


def set_title(slide, text, size=32):
    t = slide.shapes.title
    t.text = nb(text)
    for p in t.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)


def bullets(placeholder, items, l0=19, l1=16):
    """items: list of (text, level, bold)."""
    tf = placeholder.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, it in enumerate(items):
        text = it[0]
        lvl = it[1] if len(it) > 1 else 0
        bold = it[2] if len(it) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(8)
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = nb(text)
        r.font.size = Pt(l0 if lvl == 0 else l1)
        r.font.bold = bold
        r.font.color.rgb = SLATE


def pic_fit(slide, path, bx, by, bw, bh, halign="center", valign="middle"):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    if bw / bh > ar:
        h, w = bh, int(bh * ar)
    else:
        w, h = bw, int(bw / ar)
    x = {"center": bx + (bw - w) // 2, "left": bx, "right": bx + (bw - w)}[halign]
    y = {"middle": by + (bh - h) // 2, "top": by, "bottom": by + (bh - h)}[valign]
    return slide.shapes.add_picture(path, int(x), int(y), int(w), int(h))


def textbox(slide, x, y, w, h, lines, size=13, color=SLATE, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 36000
    tf.margin_top = tf.margin_bottom = 18000
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = nb(ln)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def finalize(slide, num):
    """KU Leuven footer bar: course footer (right) and slide number (left)."""
    textbox(slide, 6033600, 6210000, 4993200, 648000, [FOOTER],
            size=10, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, 576000, 6210000, 648000, 648000, [str(num)],
            size=10, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add(name):
    return prs.slides.add_slide(layout(name))


delete_all_slides()
TWO = "Twee objecten"
ONE = "Titel en object"

# ====================================================================
# SLIDE 1 - title
# ====================================================================
s = add("Titeldia")
t = s.shapes.title.text_frame
r = t.paragraphs[0].add_run()
r.text = "Multiresolution topology optimization"
r.font.size = Pt(40)
r.font.bold = True
p = t.add_paragraph()
p.space_before = Pt(12)
r = p.add_run()
r.text = ("Computational-efficiency assessment for the "
          + " ".join(["half", "MBB", "beam"]))
r.font.size = Pt(20)
sub = ph(s, 1)
sub.left, sub.top, sub.width, sub.height = MARGIN, 5230000, 6380000, 1250000
stf = sub.text_frame
stf.word_wrap = True
sublines = [("Lukas Campaert  ·  Tiebert Lefebure", 16, True),
            ("Structural Optimization  ·  B-KUL-H0T88A  ·  Prof. M. Schevenels", 13, False),
            ("Faculty of Engineering Science  ·  Academic Year 2025-2026", 13, False)]
for i, (txt, sz, bd) in enumerate(sublines):
    p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = bd
pp = ph(s, 10)
if pp is not None:
    box = (pp.left, pp.top, pp.width, pp.height)
    drop(pp)
    pic_fit(s, os.path.join(FIG, "mtop_oc_sensitivity_design.png"), *box)
notes(s, (
    "PRESENTER: Lukas   |   Target ~20 s\n\n"
    "Greeting. Our examination assignment is on multiresolution topology "
    "optimization, MTOP. In topology optimization the question is where to "
    "put material to make a structure as stiff as possible; here the question "
    "is whether we can answer that question faster. We assessed MTOP on the "
    "half MBB beam and compared it head-to-head with the classical approach. "
    "The image is our optimized beam."))

# ====================================================================
# SLIDE 2 - problem statement
# ====================================================================
s = add(TWO)
set_title(s, "Problem statement and research question")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("Minimum-compliance topology optimization of the half MBB beam", 0),
    ("MTOP (Nguyen et al., 2010): describe the design on a fine density "
     "mesh, solve the equilibrium on a coarser mesh", 0),
    ("A smaller system K u = f should lower the computation time", 0),
    ("The original paper reports no computation times or speedup factors", 0),
    ("Research question: how large is the speedup, and is the optimized "
     "design preserved?", 0, True),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(GEN, "gen_idea.png"), RIGHT_X, CONTENT_T, RIGHT_W, CONTENT_H)
finalize(s, 2)
notes(s, (
    "PRESENTER: Lukas   |   Target ~50 s\n\n"
    "Topology optimization distributes a fixed amount of material so the "
    "structure is as stiff as possible. The cost is dominated by repeatedly "
    "solving the finite element system K u = f.\n"
    "MTOP, introduced by Nguyen and co-workers in 2010, splits the problem: "
    "the design is still described on a fine mesh, but the equilibrium "
    "equations are solved on a coarse mesh, which should be much cheaper.\n"
    "The catch: the original paper never reported computation times. So our "
    "research question is exactly that - how large is the speedup, and does "
    "the optimized design stay the same?"))

# ====================================================================
# SLIDE 3 - structural model
# ====================================================================
s = add(TWO)
set_title(s, "The structural model: half MBB beam")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("Half MBB beam, plane stress, bilinear Q4 elements", 0),
    ("Unit point load at the upper-left corner", 0),
    ("Symmetry support on the left edge, roller at the lower-right corner", 0),
    ("Volume fraction 0.5, SIMP penalization power 3", 0),
    ("Cone filter radius of 24 density cells", 0),
    ("These settings are fixed for every experiment", 0, True),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(RFIG, "geometry.png"), RIGHT_X, CONTENT_T, RIGHT_W, CONTENT_H)
finalize(s, 3)
notes(s, (
    "PRESENTER: Lukas   |   Target ~40 s\n\n"
    "Our test structure is the half MBB beam from the Chapter 9 lecture code. "
    "A unit point load at the top-left corner, a symmetry condition on the "
    "left edge, and a roller at the bottom-right. Plane stress, bilinear quad "
    "elements.\n"
    "The optimization targets a volume fraction of one half, SIMP "
    "penalization three, and a cone filter of radius 24 density cells. "
    "Important: these settings are identical for every experiment, so any "
    "difference we see is due to MTOP, not the parameters."))

# ====================================================================
# SLIDE 4 - MTOP concept
# ====================================================================
s = add(TWO)
set_title(s, "The MTOP idea: decouple analysis from design")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("Classical: analysis mesh = density mesh (600 × 200)", 0),
    ("MTOP: coarse 120 × 40 analysis mesh, 5 × 5 density cells per element", 0),
    ("Same 600 × 200 design resolution, far fewer equations", 0),
    ("Element stiffness by midpoint quadrature; 25 sub-cell templates "
     "precomputed once per run", 0),
    ("Free degrees of freedom: 241,000 → 9,900", 0, True),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(GEN, "gen_concept_meshes.png"),
        RIGHT_X - 250000, CONTENT_T, RIGHT_W + 250000, CONTENT_H)
finalize(s, 4)
notes(s, (
    "PRESENTER: Lukas   |   Target ~70 s\n\n"
    "This is the heart of MTOP. In the classical approach the analysis mesh "
    "and the design mesh are the same, 600 by 200.\n"
    "In MTOP we keep the 600 by 200 design resolution, but we analyze on a "
    "coarse 120 by 40 mesh, with 5 by 5 density cells inside each finite "
    "element - that is the zoom on the right.\n"
    "Each element's stiffness is integrated with one midpoint sample per "
    "density cell, so 25 small stiffness templates, which are precomputed "
    "once per run.\n"
    "The payoff is the number of equations: the free degrees of freedom drop "
    "from about 241,000 to about 9,900. That factor is the speedup we are "
    "chasing."))

# ====================================================================
# SLIDE 5 - optimization problem + experiment matrix
# ====================================================================
s = add(TWO)
set_title(s, "Optimization problem and experiment matrix")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("Design variables: the density of each density cell", 0),
    ("Objective: minimize compliance (maximize stiffness)", 0),
    ("Constraint: material volume fraction at most 0.5", 0),
    ("SIMP material interpolation, penalization power 3", 0),
    ("Cone filter, radius 24 cells (sensitivity or density)", 0),
    ("Seven experiments: OC versus MMA, three filter / projection choices", 0, True),
])
drop(ph(s, 13))
rows = [
    ("Case", "Approach", "Optimizer", "Filter / projection"),
    ("1", "Classical", "OC", "Sensitivity filter"),
    ("2", "MTOP", "OC", "Sensitivity filter"),
    ("3", "Classical", "OC", "Density filter"),
    ("4", "MTOP", "OC", "Density filter"),
    ("5", "MTOP", "MMA", "Sensitivity filter"),
    ("6", "MTOP", "MMA", "Density filter"),
    ("7", "MTOP", "MMA", "Heaviside η = 0.3 / 0.5 / 0.7"),
]
tbl_h = 3500000
tbl_y = CONTENT_T + (CONTENT_H - tbl_h) // 2
gframe = s.shapes.add_table(len(rows), 4, RIGHT_X, tbl_y, RIGHT_W, tbl_h)
table = gframe.table
table.first_row = False
table.horz_banding = False
table.columns[0].width = 720000
table.columns[1].width = 1500000
table.columns[2].width = 1180000
table.columns[3].width = RIGHT_W - 720000 - 1500000 - 1180000
for ri, row in enumerate(rows):
    table.rows[ri].height = tbl_h // len(rows)
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = 90000
        cell.margin_right = 60000
        cell.margin_top = cell.margin_bottom = 20000
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL if ri == 0 else WHITE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if ci != 3 else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = val
        run.font.size = Pt(12.5 if ri == 0 else 12)
        run.font.bold = (ri == 0)
        run.font.color.rgb = WHITE if ri == 0 else SLATE
finalize(s, 5)
notes(s, (
    "PRESENTER: Lukas   |   Target ~50 s\n\n"
    "Quickly, the optimization problem. The design variables are the cell "
    "densities, the objective is to minimize compliance, and the constraint "
    "is the volume fraction of one half. SIMP with penalization three, and a "
    "cone filter of radius 24, used either as a sensitivity filter or a "
    "density filter.\n"
    "We ran the seven experiments in the table: the optimality-criteria runs "
    "for assignment steps 1 to 3, then MMA for step 4 with sensitivity "
    "filtering, density filtering, and Heaviside projection at three "
    "thresholds. Every MTOP run has a classical counterpart for comparison."))

# ====================================================================
# SLIDE 6 - sensitivity verification (figure on top, bullets below)
# ====================================================================
s = add(ONE)
set_title(s, "Verifying the implementation")
fd_h = 3950000
pic_fit(s, os.path.join(FIG, "sensitivity_verification_fd_strategies.png"),
        MARGIN, CONTENT_T, SW - 2 * MARGIN, fd_h, valign="top")
tx = ph(s, 1)
tx.left, tx.top = MARGIN, CONTENT_T + fd_h + 150000
tx.width, tx.height = SW - 2 * MARGIN, CONTENT_B - (CONTENT_T + fd_h + 150000)
bullets(tx, [
    ("Verified three derivative chains with three finite-difference "
     "strategies", 0),
    ("Taylor remainder slope ≈ 2.00 confirms the analytical sensitivities", 0),
], l0=17)
finalize(s, 6)
notes(s, (
    "PRESENTER: Lukas   |   Target ~40 s\n\n"
    "Before trusting any result we verified the sensitivities - the "
    "gradients the optimizer relies on. We checked three derivative chains: "
    "the classical SIMP gradient, the MTOP per-cell gradient, and the full "
    "Heaviside chain, each with three finite-difference strategies.\n"
    "The strongest is the Taylor remainder test, the right-hand panel: it "
    "should decay quadratically with the step size, and the fitted slope is "
    "essentially 2.00 for all three chains, with errors down to ten to the "
    "minus six. So the analytical gradients are correct."))

# ====================================================================
# SLIDE 7 - OC + sensitivity filter: designs match
# ====================================================================
s = add(TWO)
set_title(s, "Steps 1 and 2: MTOP reproduces the classical design")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("Optimality criteria with sensitivity filtering", 0),
    ("Same truss load path: top chord, bottom chord, triangulated web", 0),
    ("Fine-mesh compliance: classical 224.6, MTOP 224.7 (+0.045 %)", 0),
    ("Difference is a tiny boundary-level redistribution (RMS 0.003)", 0),
    ("MTOP does not change the optimized design", 0, True),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(FIG, "mtop_oc_sensitivity_design_difference.png"),
        RIGHT_X, CONTENT_T, RIGHT_W, CONTENT_H)
finalize(s, 7)
notes(s, (
    "PRESENTER: Lukas   |   Target ~45 s\n\n"
    "Now the results. Steps 1 and 2: optimality criteria with sensitivity "
    "filtering. On the right, the classical design on top, the MTOP design "
    "in the middle, and their difference at the bottom.\n"
    "The two layouts are the same truss - a top chord, a bottom chord, a "
    "triangulated web. The difference plot is almost empty, just small "
    "boundary shifts, an RMS of three thousandths. Re-evaluated on the same "
    "fine mesh the compliances differ by 0.045 percent.\n"
    "So MTOP does not change the design. I will hand over to Tiebert for how "
    "much faster it is."))

# ====================================================================
# SLIDE 8 - OC + sensitivity filter: 21x faster
# ====================================================================
s = add(TWO)
set_title(s, "Steps 1 and 2: MTOP is 21× faster")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("Same convergence path: 94 versus 95 iterations", 0),
    ("Wall-clock time: 263 s → 12 s", 0),
    ("Speedup factor of 21 at matched compliance", 0),
    ("The MTOP curve collapses against the time axis", 0),
    ("Same design, same iterations, far less time", 0, True),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(FIG, "mtop_oc_sensitivity_convergence_time_compare.png"),
        RIGHT_X, CONTENT_T, RIGHT_W, CONTENT_H)
finalize(s, 8)
notes(s, (
    "PRESENTER: Tiebert   (takes over from Lukas)   |   Target ~40 s\n\n"
    "Thanks. So the designs are identical, but the timing is not. This plot "
    "is compliance against wall-clock time - classical in blue, MTOP in "
    "orange.\n"
    "Both runs take essentially the same number of iterations, 94 and 95, so "
    "the optimization path is unchanged. But the classical run takes 263 "
    "seconds and the MTOP run takes 12 - a factor of 21. You can see the "
    "orange curve squeezed against the time axis.\n"
    "Same design, same iteration count, 21 times faster."))

# ====================================================================
# SLIDE 9 - OC + density filter
# ====================================================================
s = add(TWO)
set_title(s, "Step 3: density filtering")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("Same comparison with a density filter instead of a sensitivity filter", 0),
    ("Designs again almost identical: 0.04 % compliance difference", 0),
    ("Speedup drops to 5×: density-filter work still runs on the fine mesh", 0),
    ("OC plus density filter never meets the design-change tolerance", 0),
    ("Stopped instead on a compliance plateau", 0, True),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(FIG, "mtop_oc_density_design_difference.png"),
        RIGHT_X, CONTENT_T, RIGHT_W, CONTENT_H)
finalize(s, 9)
notes(s, (
    "PRESENTER: Tiebert   |   Target ~50 s\n\n"
    "Step 3 repeats the comparison with a density filter instead of a "
    "sensitivity filter. The designs are again almost identical, 0.04 "
    "percent in compliance.\n"
    "But the speedup drops from 21 to about 5. The reason: the density "
    "filter adds a chain-rule convolution and a volume check that still run "
    "on the full 600 by 200 density mesh, so they do not shrink with the "
    "analysis mesh.\n"
    "One more observation: with a density filter the OC update never quite "
    "meets the design-change tolerance - it oscillates at the boundary - so "
    "we stop it on a compliance plateau instead."))

# ====================================================================
# SLIDE 10 - MMA + Heaviside
# ====================================================================
s = add(TWO)
set_title(s, "Step 4: MMA and Heaviside projection")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("OC update replaced by the method of moving asymptotes (MMA)", 0),
    ("MMA alone: design close to OC; density-filter convergence issue remains", 0),
    ("Heaviside projection pushes the design towards black-and-white", 0),
    ("Near-binary layout: ~17 % lower compliance than the gray design", 0),
    ("All three thresholds agree within 0.7 %; 0.5 is best", 0, True),
])
drop(ph(s, 13))
heavi = [("η = 0.3", "mtop_mma_heaviside_eta_030_design.png"),
         ("η = 0.5   (best)", "mtop_mma_heaviside_eta_050_design.png"),
         ("η = 0.7", "mtop_mma_heaviside_eta_070_design.png")]
slot = CONTENT_H // 3
for i, (lab, fn) in enumerate(heavi):
    sy = CONTENT_T + i * slot
    textbox(s, RIGHT_X, sy, RIGHT_W, 300000, [lab], size=13, bold=True,
            color=TEAL, align=PP_ALIGN.CENTER)
    pic_fit(s, os.path.join(FIG, fn), RIGHT_X, sy + 300000,
            RIGHT_W, slot - 360000, valign="top")
finalize(s, 10)
notes(s, (
    "PRESENTER: Tiebert   |   Target ~65 s\n\n"
    "Step 4 swaps the optimality-criteria update for MMA, the method of "
    "moving asymptotes. MMA on its own gives a design close to the OC "
    "result, and it does not by itself fix the density-filter convergence "
    "issue from the previous slide.\n"
    "The real benefit comes with Heaviside projection, which pushes the gray "
    "material towards either solid or void. On the right are the three "
    "projected designs, for thresholds 0.3, 0.5 and 0.7: all near "
    "black-and-white, all the same load path.\n"
    "The projected compliance is about 17 percent below the gray "
    "density-filtered design. The three thresholds differ by less than 0.7 "
    "percent; 0.5 is marginally the best."))

# ====================================================================
# SLIDE 11 - where the speedup comes from (figure left, text right)
# ====================================================================
s = add(TWO)
set_title(s, "Where does the speedup come from?")
fig_w = 6700000
pic_fit(s, os.path.join(FIG, "phase_breakdown_chart.png"),
        MARGIN, CONTENT_T, fig_w, CONTENT_H, halign="left")
tx = ph(s, 1)
tx.left = MARGIN + fig_w + 280000
tx.top = CONTENT_T
tx.width = SW - MARGIN - tx.left
tx.height = CONTENT_H
bullets(tx, [
    ("Classical OC: 82 % of the time is the finite element solve", 0),
    ("MTOP shrinks that solve by about 40×", 0),
    ("Filtering, assembly and the optimizer still act on the full density "
     "mesh, so they cap the speedup", 0),
    ("MMA runs become dominated by the optimizer subproblem", 0),
], l0=17)
drop(ph(s, 13))
finalize(s, 11)
notes(s, (
    "PRESENTER: Tiebert   |   Target ~50 s\n\n"
    "So where does the speedup actually come from? This chart breaks each "
    "run into algorithmic phases.\n"
    "For the classical OC run, the green block - the finite element solve - "
    "is 82 percent of the time. MTOP shrinks exactly that block: the coarse "
    "solve is about 40 times faster.\n"
    "What it cannot shrink is the rest - the filtering, the assembly, the "
    "optimizer - which still act on the fine density mesh. That is why the "
    "total speedup is 21 and not 40. For the MMA runs the red optimizer "
    "block dominates, so there the equilibrium solve is no longer the "
    "bottleneck."))

# ====================================================================
# SLIDE 12 - what MTOP really buys you
# ====================================================================
s = add(TWO)
set_title(s, "What MTOP really buys you")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("A plain coarse model reaches within 0.6 % of the fine design", 0),
    ("For this 2D case with a large filter, a coarse model already captures "
     "the load path", 0),
    ("But that coarse design is grayer and fixed at the coarse resolution", 0),
    ("MTOP keeps the full 600 × 200 design while paying only the coarse "
     "solve", 0, True),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(FIG, "classical_oc_coarse_sensitivity_design.png"),
        RIGHT_X, CONTENT_T, RIGHT_W, CONTENT_H - 300000)
textbox(s, RIGHT_X, CONTENT_B - 300000, RIGHT_W, 300000,
        ["Coarse-resolution baseline: 120 × 40 analysis = density mesh"],
        size=12, color=MUTED, align=PP_ALIGN.CENTER)
finalize(s, 12)
notes(s, (
    "PRESENTER: Tiebert   |   Target ~50 s\n\n"
    "One honest qualification. We also ran a plain coarse model, 120 by 40 "
    "for both the analysis and the design. Re-analyzed on the fine mesh, it "
    "lands within 0.6 percent of the full design.\n"
    "So for this 2D beam, with a fairly large filter, even a crude model "
    "already finds the load path. The difference is that this coarse design "
    "is locked to 120-by-40 resolution and is noticeably grayer.\n"
    "MTOP's real value is that it keeps the full 600-by-200 design "
    "description while only paying for the coarse solve."))

# ====================================================================
# SLIDE 13 - conclusions
# ====================================================================
s = add(TWO)
set_title(s, "Conclusions")
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = LEFT_X, CONTENT_T, LEFT_W, CONTENT_H
bullets(tx, [
    ("MTOP reproduces the classical MBB design to within 0.05 % compliance", 0),
    ("Sensitivity filtering: 21× faster; density filtering: 5× faster", 0),
    ("The speedup comes from the cheaper coarse equilibrium solve", 0),
    ("OC plus density filtering converges only in the compliance sense", 0),
    ("Heaviside projection gives the lowest, near-binary compliance "
     "(threshold 0.5)", 0),
])
drop(ph(s, 13))
pic_fit(s, os.path.join(FIG, "runtime_summary_chart.png"),
        RIGHT_X, CONTENT_T, RIGHT_W, CONTENT_H)
finalize(s, 13)
notes(s, (
    "PRESENTER: Tiebert   |   Target ~40 s\n\n"
    "To conclude. MTOP reproduces the classical MBB design to within 0.05 "
    "percent compliance. It is 21 times faster with sensitivity filtering "
    "and 5 times faster with density filtering.\n"
    "The speedup comes from the cheaper coarse equilibrium solve, and is "
    "capped by the operations that still run on the fine mesh. OC with "
    "density filtering converges only in the compliance sense. And Heaviside "
    "projection gives the lowest, near-binary compliance, best at threshold "
    "0.5.\n"
    "The chart summarizes all nine runs - cost on the left, compliance on "
    "the right."))

# ====================================================================
# SLIDE 14 - use of GenAI
# ====================================================================
s = add(TWO)
set_title(s, "Use of GenAI")
col_w = (SW - 2 * MARGIN - COLGAP) // 2
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = MARGIN, CONTENT_T, col_w, CONTENT_H
bullets(tx, [
    ("We used ChatGPT for", 0, True),
    ("Language editing of text we wrote ourselves", 1),
    ("Post-processing and plotting scripts", 1),
    ("One supporting theory figure (the MMA sketch)", 1),
], l0=21, l1=18)
rx = ph(s, 13)
rx.left, rx.top, rx.width, rx.height = MARGIN + col_w + COLGAP, CONTENT_T, col_w, CONTENT_H
bullets(rx, [
    ("We kept the work in our control", 0, True),
    ("MTOP method and MATLAB solver implemented by us", 1),
    ("All sensitivities verified with finite differences", 1),
    ("Every result reproduced and understood ourselves", 1),
], l0=21, l1=18)
finalize(s, 14)
notes(s, (
    "PRESENTER: Tiebert   |   Target ~20 s\n\n"
    "A brief word on GenAI. We used ChatGPT as a support tool: to "
    "language-edit text we had written ourselves, for post-processing and "
    "plotting scripts, and for one supporting theory figure, the MMA "
    "sketch.\n"
    "We did not use it for the method or the results. The MTOP "
    "implementation is ours, every sensitivity was verified independently "
    "with finite differences, and every number in this talk we reproduced "
    "and understand. The full transparency table is in the report."))

# ====================================================================
# SLIDE 15 - future work and outlook
# ====================================================================
s = add(TWO)
set_title(s, "Future work and outlook")
col_w = (SW - 2 * MARGIN - COLGAP) // 2
col_h = 3500000
tx = ph(s, 1)
tx.left, tx.top, tx.width, tx.height = MARGIN, CONTENT_T, col_w, col_h
bullets(tx, [
    ("Where MTOP pays off most", 0, True),
    ("3D problems: the equilibrium solve dominates far more, so the "
     "speedup is larger", 1),
    ("Smaller filter radii and finer features, which a plain coarse model "
     "cannot resolve", 1),
    ("Cases that need the fine design description, not only the overall "
     "load path", 1),
], l0=20, l1=17)
rx = ph(s, 13)
rx.left, rx.top, rx.width, rx.height = MARGIN + col_w + COLGAP, CONTENT_T, col_w, col_h
bullets(rx, [
    ("Open issues and limitations", 0, True),
    ("Density filtering needs a convergence criterion beyond the "
     "design-change tolerance", 1),
    ("Filtering, assembly and the optimizer still scale with the fine "
     "mesh and cap the speedup", 1),
    ("With MMA the optimizer subproblem, not the FE solve, is the "
     "bottleneck", 1),
], l0=20, l1=17)
textbox(s, MARGIN, CONTENT_T + col_h + 200000, SW - 2 * MARGIN, 680000,
        ["Thank you for your attention. Questions?"],
        size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
finalize(s, 15)
notes(s, (
    "PRESENTER: Tiebert   |   Target ~30 s\n\n"
    "To close, a brief outlook. Our headline numbers are for a 2D beam, "
    "where a plain coarse model already came within 0.6 percent, so the 2D "
    "benefit was modest. MTOP pays off most where the equilibrium solve "
    "genuinely dominates the cost: in 3D, with much larger systems, and "
    "where the design has fine features a coarse model cannot resolve, so "
    "keeping the fine density description is worth it.\n"
    "Three open issues from this study: density filtering needs a better "
    "convergence criterion than the design-change tolerance; the filtering, "
    "assembly and optimizer steps still scale with the fine mesh and cap the "
    "speedup below the 40 times of the solve itself; and with MMA the "
    "optimizer subproblem, not the FE solve, becomes the bottleneck.\n"
    "Thank you - we are happy to take questions."))

prs.save(OUT)
print("saved", OUT, "with", len(prs.slides._sldIdLst), "slides")
